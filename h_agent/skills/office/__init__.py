"""
h_agent/skills/office - Windows Office Automation Skill

Provides automation for Microsoft Office applications:
- Word: Document creation, editing, formatting
- Excel: Spreadsheet operations, formulas, charts
- PowerPoint: Presentation creation and editing

Dependencies (Windows only):
    pip install python-docx openpyxl python-pptx pywin32

Usage:
    from h_agent.skills.office import Word, Excel, PowerPoint
    
    # Word
    doc = Word.create_document("output.docx")
    Word.add_heading(doc, "Title")
    Word.add_paragraph(doc, "Content")
    Word.save(doc)
    
    # Excel
    wb = Excel.create_workbook("output.xlsx")
    Excel.write_cell(wb, "Sheet1", "A1", "Value")
    Excel.save(wb)
    
    # PowerPoint
    ppt = PowerPoint.create_presentation("output.pptx")
    PowerPoint.add_slide(ppt, "Title Slide")
    PowerPoint.save(ppt)
"""

import os
import sys
import platform
from pathlib import Path
from typing import Optional, List, Dict, Any, Union

# Skill metadata
SKILL_NAME = "office"
SKILL_VERSION = "1.0.0"
SKILL_DESCRIPTION = "Windows Office automation (Word, Excel, PowerPoint)"
SKILL_AUTHOR = "h-agent team"
SKILL_CATEGORY = "office"
SKILL_DEPENDENCIES = ["docx", "openpyxl", "pptx"]
SKILL_PLATFORMS = ["windows"]
SKILL_TOOLS = []
SKILL_FUNCTIONS = {}

# Check platform on import
if platform.system().lower() != "windows":
    def _windows_only():
        raise OSError("Office skill is only available on Windows")
    _windows_only_msg = _windows_only
else:
    _windows_only_msg = None


def _check_dependencies():
    """Check if required dependencies are installed."""
    missing = []
    for dep in SKILL_DEPENDENCIES:
        try:
            __import__(dep)
        except ImportError:
            missing.append(dep)
    if missing:
        raise ImportError(
            f"Missing dependencies for office skill: {', '.join(missing)}\n"
            f"Install with: pip install {' '.join(SKILL_DEPENDENCIES)}"
        )


# ─────────────────────────────────────────────
# Word Operations
# ─────────────────────────────────────────────

class Word:
    """Microsoft Word automation."""
    
    @staticmethod
    def create_document(path: Optional[str] = None) -> Any:
        """
        Create a new Word document.
        
        Args:
            path: Output file path (optional)
            
        Returns:
            Document object
        """
        _check_dependencies()
        from docx import Document
        return Document()
    
    @staticmethod
    def open_document(path: str) -> Any:
        """
        Open an existing Word document.
        
        Args:
            path: File path to open
            
        Returns:
            Document object
        """
        _check_dependencies()
        from docx import Document
        if not os.path.exists(path):
            raise FileNotFoundError(f"Document not found: {path}")
        return Document(path)
    
    @staticmethod
    def add_heading(doc: Any, text: str, level: int = 1) -> Any:
        """Add a heading to the document."""
        return doc.add_heading(text, level=level)
    
    @staticmethod
    def add_paragraph(doc: Any, text: str, style: Optional[str] = None) -> Any:
        """Add a paragraph to the document."""
        return doc.add_paragraph(text, style=style)
    
    @staticmethod
    def add_table(doc: Any, rows: int, cols: int, data: Optional[List[List[str]]] = None) -> Any:
        """Add a table to the document."""
        table = doc.add_table(rows=rows, cols=cols)
        if data:
            for i, row_data in enumerate(data):
                row = table.rows[i]
                for j, cell_data in enumerate(row_data):
                    row.cells[j].text = str(cell_data)
        return table
    
    @staticmethod
    def write_cell(table: Any, row: int, col: int, value: str) -> None:
        """Write to a table cell."""
        table.rows[row].cells[col].text = str(value)
    
    @staticmethod
    def add_image(doc: Any, image_path: str, width: Optional[float] = None) -> Any:
        """Add an image to the document."""
        from docx.shared import Inches
        if width:
            return doc.add_picture(image_path, width=Inches(width))
        return doc.add_picture(image_path)
    
    @staticmethod
    def set_page_margins(doc: Any, top: float = 1.0, bottom: float = 1.0, 
                         left: float = 1.0, right: float = 1.0) -> None:
        """Set page margins in inches."""
        from docx.shared import Inches
        section = doc.sections[0]
        section.top_margin = Inches(top)
        section.bottom_margin = Inches(bottom)
        section.left_margin = Inches(left)
        section.right_margin = Inches(right)
    
    @staticmethod
    def save(doc: Any, path: Optional[str] = None) -> str:
        """
        Save the document.
        
        Args:
            doc: Document object
            path: Output path (required if not opening existing)
            
        Returns:
            Saved file path
        """
        if not path:
            raise ValueError("Path required to save document")
        doc.save(path)
        return path
    
    @staticmethod
    def get_text(doc: Any) -> str:
        """Extract all text from document."""
        return "\n".join([p.text for p in doc.paragraphs])


# ─────────────────────────────────────────────
# Excel Operations
# ─────────────────────────────────────────────

class Excel:
    """Microsoft Excel automation."""
    
    @staticmethod
    def create_workbook(path: Optional[str] = None, sheet_name: str = "Sheet1") -> Any:
        """
        Create a new Excel workbook.
        
        Args:
            path: Output file path (optional)
            sheet_name: Name of the first sheet
            
        Returns:
            Workbook object
        """
        _check_dependencies()
        from openpyxl import Workbook
        wb = Workbook()
        if sheet_name != "Sheet":
            ws = wb.active
            ws.title = sheet_name
        return wb
    
    @staticmethod
    def open_workbook(path: str) -> Any:
        """
        Open an existing Excel workbook.
        
        Args:
            path: File path to open
            
        Returns:
            Workbook object
        """
        _check_dependencies()
        from openpyxl import load_workbook
        if not os.path.exists(path):
            raise FileNotFoundError(f"Workbook not found: {path}")
        return load_workbook(path)
    
    @staticmethod
    def get_sheet(wb: Any, name: str) -> Any:
        """Get a sheet by name."""
        if name not in wb.sheetnames:
            raise ValueError(f"Sheet not found: {name}")
        return wb[name]
    
    @staticmethod
    def create_sheet(wb: Any, name: str, index: Optional[int] = None) -> Any:
        """Create a new sheet."""
        return wb.create_sheet(name, index=index)
    
    @staticmethod
    def delete_sheet(wb: Any, name: str) -> None:
        """Delete a sheet."""
        if name in wb.sheetnames:
            del wb[name]
    
    @staticmethod
    def write_cell(wb: Any, sheet: str, cell: str, value: Any) -> None:
        """
        Write a value to a cell.
        
        Args:
            wb: Workbook object
            sheet: Sheet name
            cell: Cell reference (e.g., "A1")
            value: Value to write
        """
        ws = wb[sheet] if isinstance(sheet, str) else sheet
        ws[cell] = value
    
    @staticmethod
    def read_cell(wb: Any, sheet: str, cell: str) -> Any:
        """Read a value from a cell."""
        ws = wb[sheet] if isinstance(sheet, str) else sheet
        return ws[cell].value
    
    @staticmethod
    def write_range(wb: Any, sheet: str, start_cell: str, data: List[List[Any]]) -> None:
        """
        Write a 2D array to a range.
        
        Args:
            wb: Workbook object
            sheet: Sheet name
            start_cell: Starting cell (e.g., "A1")
            data: 2D list of values
        """
        from openpyxl.utils import get_column_letter
        ws = wb[sheet] if isinstance(sheet, str) else sheet
        
        # Parse start cell
        col_str = "".join(c for c in start_cell if c.isalpha())
        row_start = int("".join(c for c in start_cell if c.isdigit()))
        col_start = sum((ord(c.upper()) - ord('A') + 1) * (26 ** i) 
                        for i, c in enumerate(reversed(col_str.upper())))
        
        for i, row in enumerate(data):
            for j, value in enumerate(row):
                col = get_column_letter(col_start + j)
                ws[f"{col}{row_start + i}"] = value
    
    @staticmethod
    def set_column_width(ws: Any, column: str, width: float) -> None:
        """Set column width."""
        ws.column_dimensions[column].width = width
    
    @staticmethod
    def set_row_height(ws: Any, row: int, height: float) -> None:
        """Set row height."""
        ws.row_dimensions[row].height = height
    
    @staticmethod
    def merge_cells(ws: Any, start_cell: str, end_cell: str) -> None:
        """Merge cells."""
        ws.merge_cells(f"{start_cell}:{end_cell}")
    
    @staticmethod
    def add_formula(ws: Any, cell: str, formula: str) -> None:
        """Add a formula to a cell."""
        ws[cell] = formula
    
    @staticmethod
    def apply_style(ws: Any, cell: str, bold: bool = False, italic: bool = False,
                    font_size: Optional[int] = None, font_color: Optional[str] = None) -> None:
        """Apply text style to a cell."""
        from openpyxl.styles import Font
        cell_obj = ws[cell]
        current = cell_obj.font
        cell_obj.font = Font(
            name=current.name,
            size=font_size or current.size,
            bold=bold,
            italic=italic,
            color=font_color or current.color
        )
    
    @staticmethod
    def save(wb: Any, path: Optional[str] = None) -> str:
        """
        Save the workbook.
        
        Args:
            wb: Workbook object
            path: Output path (required for new workbooks)
            
        Returns:
            Saved file path
        """
        if not path:
            raise ValueError("Path required to save workbook")
        wb.save(path)
        return path
    
    @staticmethod
    def get_sheet_names(wb: Any) -> List[str]:
        """Get all sheet names."""
        return wb.sheetnames
    
    @staticmethod
    def get_data_range(ws: Any) -> List[List[Any]]:
        """Get all data in sheet as 2D list."""
        data = []
        for row in ws.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                data.append(list(row))
        return data


# ─────────────────────────────────────────────
# PowerPoint Operations
# ─────────────────────────────────────────────

class PowerPoint:
    """Microsoft PowerPoint automation."""
    
    @staticmethod
    def create_presentation(path: Optional[str] = None) -> Any:
        """
        Create a new PowerPoint presentation.
        
        Args:
            path: Output file path (optional)
            
        Returns:
            Presentation object
        """
        _check_dependencies()
        from pptx import Presentation
        return Presentation()
    
    @staticmethod
    def open_presentation(path: str) -> Any:
        """
        Open an existing PowerPoint presentation.
        
        Args:
            path: File path to open
            
        Returns:
            Presentation object
        """
        _check_dependencies()
        from pptx import Presentation
        if not os.path.exists(path):
            raise FileNotFoundError(f"Presentation not found: {path}")
        return Presentation(path)
    
    @staticmethod
    def add_slide(prs: Any, layout_title: str = "Title Slide") -> Any:
        """
        Add a slide to the presentation.
        
        Args:
            prs: Presentation object
            layout_title: Layout name to use
            
        Returns:
            Slide object
        """
        # Find layout by title
        layout = None
        for l in prs.slide_layouts:
            if l.name.lower() == layout_title.lower():
                layout = l
                break
        
        if layout is None:
            layout = prs.slide_layouts[0]  # Default to first layout
        
        return prs.slides.add_slide(layout)
    
    @staticmethod
    def set_slide_title(slide: Any, title: str) -> None:
        """Set the title of a slide."""
        if slide.shapes.title:
            slide.shapes.title.text = title
    
    @staticmethod
    def add_text(slide: Any, text: str, left: float = 0.5, top: float = 1.5,
                 width: float = 8.0, height: float = 1.0) -> Any:
        """
        Add a text box to a slide.
        
        Args:
            slide: Slide object
            text: Text content
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            
        Returns:
            Text box shape
        """
        from pptx.util import Inches
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), 
                                           Inches(width), Inches(height))
        textbox.text_frame.text = text
        return textbox
    
    @staticmethod
    def add_image(slide: Any, image_path: str, left: float = 1.0, top: float = 1.0,
                  width: Optional[float] = None) -> Any:
        """
        Add an image to a slide.
        
        Args:
            slide: Slide object
            image_path: Path to image file
            left: Left position in inches
            top: Top position in inches
            width: Width in inches (optional, maintains aspect ratio)
            
        Returns:
            Picture shape
        """
        from pptx.util import Inches
        if width:
            pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top), 
                                           width=Inches(width))
        else:
            pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top))
        return pic
    
    @staticmethod
    def add_table(slide: Any, data: List[List[str]], left: float = 0.5, 
                  top: float = 2.0, width: float = 9.0) -> Any:
        """Add a table to a slide."""
        from pptx.util import Inches
        rows = len(data)
        cols = max(len(row) for row in data) if data else 0
        
        table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                       Inches(width), Inches(1.0)).table
        
        for i, row_data in enumerate(data):
            for j, cell_data in enumerate(row_data):
                table.cell(i, j).text = str(cell_data)
        
        return table
    
    @staticmethod
    def set_background_color(prs: Any, color: str) -> None:
        """Set the background color of all slides."""
        from pptx.dml.color import RGBColor
        for slide in prs.slides:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(color.replace("#", ""))
    
    @staticmethod
    def save(prs: Any, path: Optional[str] = None) -> str:
        """
        Save the presentation.
        
        Args:
            prs: Presentation object
            path: Output path (required for new presentations)
            
        Returns:
            Saved file path
        """
        if not path:
            raise ValueError("Path required to save presentation")
        prs.save(path)
        return path
    
    @staticmethod
    def get_slide_count(prs: Any) -> int:
        """Get the number of slides."""
        return len(prs.slides)
    
    @staticmethod
    def get_slide_titles(prs: Any) -> List[str]:
        """Get titles of all slides."""
        titles = []
        for slide in prs.slides:
            if slide.shapes.title:
                titles.append(slide.shapes.title.text)
            else:
                titles.append("")
        return titles


# ─────────────────────────────────────────────
# Register skill functions
# ─────────────────────────────────────────────

SKILL_FUNCTIONS = {
    # Word
    "word_create": Word.create_document,
    "word_open": Word.open_document,
    "word_add_heading": Word.add_heading,
    "word_add_paragraph": Word.add_paragraph,
    "word_add_table": Word.add_table,
    "word_add_image": Word.add_image,
    "word_save": Word.save,
    "word_get_text": Word.get_text,
    
    # Excel
    "excel_create": Excel.create_workbook,
    "excel_open": Excel.open_workbook,
    "excel_get_sheet": Excel.get_sheet,
    "excel_create_sheet": Excel.create_sheet,
    "excel_write_cell": Excel.write_cell,
    "excel_read_cell": Excel.read_cell,
    "excel_write_range": Excel.write_range,
    "excel_add_formula": Excel.add_formula,
    "excel_save": Excel.save,
    "excel_get_sheet_names": Excel.get_sheet_names,
    "excel_get_data_range": Excel.get_data_range,
    
    # PowerPoint
    "ppt_create": PowerPoint.create_presentation,
    "ppt_open": PowerPoint.open_presentation,
    "ppt_add_slide": PowerPoint.add_slide,
    "ppt_set_slide_title": PowerPoint.set_slide_title,
    "ppt_add_text": PowerPoint.add_text,
    "ppt_add_image": PowerPoint.add_image,
    "ppt_add_table": PowerPoint.add_table,
    "ppt_save": PowerPoint.save,
    "ppt_get_slide_count": PowerPoint.get_slide_count,
    "ppt_get_slide_titles": PowerPoint.get_slide_titles,
}
